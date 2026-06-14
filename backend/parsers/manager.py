import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import os
from typing import List, Dict, Any
from loguru import logger

class ParserManager:
    @staticmethod
    def parse_pdf(file_path: str) -> List[Dict[str, Any]]:
        pages = []
        try:
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc):
                text = page.get_text()
                pages.append({
                    "text": text,
                    "metadata": {
                        "page_number": page_num + 1,
                        "source_type": "pdf"
                    }
                })
            doc.close()
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
        return pages

    @staticmethod
    def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
        slides = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                slides.append({
                    "text": text,
                    "metadata": {
                        "slide_number": i + 1,
                        "source_type": "pptx"
                    }
                })
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
        return slides

    @staticmethod
    def parse_docx(file_path: str) -> List[Dict[str, Any]]:
        try:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return [{
                "text": text,
                "metadata": {
                    "source_type": "docx"
                }
            }]
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
        return []

    @staticmethod
    def parse_txt(file_path: str) -> List[Dict[str, Any]]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            return [{
                "text": text,
                "metadata": {
                    "source_type": "txt"
                }
            }]
        except Exception as e:
            logger.error(f"Error parsing TXT {file_path}: {e}")
        return []

    def parse(self, file_path: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self.parse_pdf(file_path)
        elif ext == ".pptx":
            return self.parse_pptx(file_path)
        elif ext == ".docx":
            return self.parse_docx(file_path)
        elif ext == ".txt":
            return self.parse_txt(file_path)
        else:
            logger.warning(f"Unsupported file extension: {ext}")
            return []
